#!/usr/bin/python 
# -*- coding: <encoding name> -*-
from os import system
from os import unlink
import os
try:
    import fitz
except:
    system('pip install pymupdf')
try:
    import PyPDF2
except:
    system('pip install PyPDF2')
try:
    import gtts
except:
    system('pip install gtts')
try:
    import pytesseract
except:
    system('pip install pytesseract')
#1788351105:AAHfYx9EsreTguPB1PUkDzHZ1ZhZz5m-nuc
try:
    import PIL
except:
    system('pip install PIL')
try:
    import io
except:
    system('pip install io')
try:
    import wand
except:
    system('pip install wand')
import shutil

try:
    from pptx import Presentation
except:
    system('pip install python-pptx')
try:
    from docx import Document
except:
    system('pip install python-docx')
from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Pt, RGBColor, Inches
from docx.oxml.ns import qn
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import PyPDF2

try:
    from selenium import webdriver
except:
    system('pip install selenium')
import fitz
try:
    import telebot
except:
    system('pip install pytelegrambotapi')
from deep_translator import GoogleTranslator
import telebot
from telebot import types
import requests
from PIL import Image
import io
#import pytesseract
from wand.image import Image as wi
from pptx import Presentation
from docx import Document
import glob

#pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

print('tran')
r1 = requests.session()
token = '1601467114:AAEQhSKTf8pSbSBFcSkG-KzDPlYZONee0X8'
bot = telebot.TeleBot("1601467114:AAEQhSKTf8pSbSBFcSkG-KzDPlYZONee0X8")
open('tran.txt', 'a').write('')
c = []
g = []
d = []
work = []
idi = 209501902


@bot.message_handler(commands=['start'])
def key(msg):
    ch = msg.chat.id
    try:
        if ch in c:
            pass
        else:
            c.append(ch)
        if msg.chat.id == idi:
            admin = types.ReplyKeyboardMarkup()
            akoky = types.KeyboardButton('الاذاعة')
            zkoky = types.KeyboardButton('الاعضاء')
            admin.add(akoky, zkoky)
            bot.reply_to(msg, 'اهلا و سهلا بك ايها المطور {}'.format(msg.from_user.first_name), reply_markup=admin)

        else:
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']
            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
            else:
                q = types.InlineKeyboardMarkup()
                q1 = types.InlineKeyboardButton('اضغط هنا \nاذا كان هنالك مصطلح طبي تريد ترجمته', callback_data='q1')
                q.add(q1)
                a9 = types.InlineKeyboardButton("المطور 💚", url="https://t.me/Q5QQQQ")
                a10 = types.InlineKeyboardButton("قناه الاعلانات↗️", url="https://t.me/akokybot")
                q.add(a9)
                q.add(a10)
                bot.send_message(msg.chat.id, 'اهلا و سهلا بك في بالبوت الترجمة 💚', reply_markup=q)
            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
    except:
        pass


@bot.message_handler(content_types=['photo'])
def key(msg):
    ch = msg.chat.id
    try:
        if msg.chat.type != "private":
            if ch in g:
                if ch in c:
                    pass
                else:
                    c.append(ch)

                ch = msg.chat.id
                if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                    pass
                else:
                    open('tran.txt', 'a').write(f'{msg.chat.id}\n')
                    bot.send_video(ch, 'BAACAgIAAxkBAAEBMbBgCBkFA1GYku_Q1wdBAAHTqU-ItsIAAhQMAAKHfEBIKhPv8pCI94oeBA')
                    bot.send_voice(ch, 'AwACAgIAAxkBAAEBMX9gCBK56Cj7teHoOhauUWlr7poIZQACEAwAAod8QEiJYsq2h5MWXR4E')
                try:
                    file_info = bot.get_file(
                        str(msg).split("file_id': '")[-1][
                        :str(msg).split("file_id': '")[-1].find("', 'file_unique_id")])
                    downloaded_file = bot.download_file(file_info.file_path)
                    with open('xc.jpg', 'wb') as new_file:
                        new_file.write(downloaded_file)
                    im = Image.open('xc.jpg')
                    text = pytesseract.image_to_string(im)
                    url = 'https://www.arabtran.com/gtranslate/'
                    head = {
                        'accept': '*/*',
                        'accept-encoding': 'gzip, deflate, br',
                        'accept-language': 'en-US,en;q=0.9',
                        'content-length': '31',
                        'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                        'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                        'origin': 'https://www.arabtran.com',
                        'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                        'sec-fetch-dest': 'empty',
                        'sec-fetch-mode': 'cors',
                        'sec-fetch-site': 'same-origin',
                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                        'x-requested-with': 'XMLHttpRequest',
                    }
                    data = {
                        'text': text,
                        'gfrom': 'en',
                        'gto': 'ar',
                        'key': 'ABC'
                    }
                    x = r1.post(url=url, data=data, ).text
                    if ch == idi:
                        bot.send_message(ch, text)
                    bot.send_message(ch, x)
                except:
                    pass
        else:
            if ch in c:
                pass
            else:
                c.append(ch)

            ch = msg.chat.id
            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
                bot.send_video(ch, 'BAACAgIAAxkBAAEBMbBgCBkFA1GYku_Q1wdBAAHTqU-ItsIAAhQMAAKHfEBIKhPv8pCI94oeBA')
                bot.send_voice(ch, 'AwACAgIAAxkBAAEBMX9gCBK56Cj7teHoOhauUWlr7poIZQACEAwAAod8QEiJYsq2h5MWXR4E')
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']
            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
            else:
                try:
                    file_info = bot.get_file(
                        str(msg).split("file_id': '")[-1][
                        :str(msg).split("file_id': '")[-1].find("', 'file_unique_id")])
                    downloaded_file = bot.download_file(file_info.file_path)
                    with open('xc.jpg', 'wb') as new_file:
                        new_file.write(downloaded_file)
                    im = Image.open('xc.jpg')
                    text = pytesseract.image_to_string(im)
                    url = 'https://www.arabtran.com/gtranslate/'
                    head = {
                        'accept': '*/*',
                        'accept-encoding': 'gzip, deflate, br',
                        'accept-language': 'en-US,en;q=0.9',
                        'content-length': '31',
                        'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                        'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                        'origin': 'https://www.arabtran.com',
                        'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                        'sec-fetch-dest': 'empty',
                        'sec-fetch-mode': 'cors',
                        'sec-fetch-site': 'same-origin',
                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                        'x-requested-with': 'XMLHttpRequest',
                    }
                    data = {
                        'text': text,
                        'gfrom': 'en',
                        'gto': 'ar',
                        'key': 'ABC'
                    }
                    x = GoogleTranslator(source='auto', target='ar').translate(
                        text=text)
                    bot.send_message(ch, x)
                    if ch == idi or msg.from_user.username == 'moredha_2000' or msg.from_user.username == 'Fattma_1996':
                        bot.send_message(ch, text)
                        x1 = x.split('\n')
                        x2 = text.split('\n')

                        i = 0
                        x3 = ''
                        try:
                            while i < len(x1):
                                x3 += f'{x1[i]}\n{x2[i]}\n'
                                i += 1
                        except:
                            pass

                        if len(x3) < 5000:
                            bot.send_message(ch, x3[:5000])
                            bot.send_message(ch, x3[5000:])
                        else:
                            bot.send_message(ch, x3)

                except:
                    pass
        q = types.InlineKeyboardMarkup()
        a9 = types.InlineKeyboardButton("المطور 💚", url="https://t.me/Q5QQQQ")
        a10 = types.InlineKeyboardButton("قناه الاعلانات↗️", url="https://t.me/akokybot")
        q.add(a9)
        q.add(a10)
        bot.send_message(ch, 'الإعلانات ⬇,️', reply_markup=q)
    except:
        pass


@bot.message_handler(content_types='text')
def an(msg):
    ch = msg.chat.id
    x = msg.text
    try:
        if msg.chat.type != "private":
            if ch in g:
                if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                    pass
                else:
                    open('tran.txt', 'a').write(f'{msg.chat.id}\n')
                try:
                    if msg.text == 'الاعضاء' and msg.chat.id == idi:
                        x = open('tran.txt', 'r').readlines()
                        bot.reply_to(msg, '{}'.format(len(x)))
                except:
                    pass
                try:
                    if msg.text == 'الاذاعة' and msg.chat.id == idi:
                        markup = types.ForceReply(selective=False)
                        bot.send_message(msg.chat.id, "ارسل اذاعتك", reply_markup=markup)
                except:
                    pass
                try:
                    if msg.reply_to_message.text == "ارسل اذاعتك":
                        try:
                            s = msg.text
                            x = open('tran.txt', 'r')
                            for i in x:
                                try:
                                    bot.send_message(i.replace('\n', ''), s)
                                except:
                                    pass
                        except:
                            pass
                except:
                    pass
                w = msg.text.lower()
                try:
                    if msg.reply_to_message.text == "ارسل المصطلح الطبي المراد ترجمته":
                        if w[0] == 'ا' or w[0] == 'أ' and w[1] == 'ل':
                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                            head = {
                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                'Accept-Encoding': 'gzip, deflate, br',
                                'Accept-Language': 'en-US,en;q=0.9',
                                'Connection': 'keep-alive',
                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                'Host': 'context.reverso.net',
                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                'Sec-Fetch-Dest': 'document',
                                'Sec-Fetch-Mode': 'navigate',
                                'Sec-Fetch-Site': 'same-origin',
                                'Sec-Fetch-User': '?1',
                                'Upgrade-Insecure-Requests': '1',
                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                            }
                            j = r1.get(url=url, headers=head)
                            v = (j.text[
                                 j.text.find('<button class="other-content" data-other="0" data-negative="') + len(
                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                     'Other translations</button>')])[:-2]
                            c = (v.replace("-{", "").replace("}", "\n"))
                            bot.send_message(msg.chat.id, f'* {c} *')

                        else:
                            try:
                                if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n',
                                                 'o',
                                                 'p',
                                                 'q',
                                                 'r',
                                                 's', 't', 'u', 'v', 'w', 'x', 'y', 'z']:
                                    d = w
                                    v = f'https://www.tbeeb.net/%D9%82%D8%A7%D9%85%D9%88%D8%B3-%D8%B7%D8%A8%D9%8A/search.php?q={d}+&dictionary=%D8%A8%D8%AD%D8%AB'
                                    try:
                                        i = 1
                                        zz = ''
                                        while i < len(r1.get(v).text.split('"tden">')):
                                            if i > 7:
                                                pass
                                            else:
                                                x1 = r1.get(v).text.split('"tden">')[i][
                                                     :r1.get(v).text.split('"tden">')[i].find('</td>')]
                                                x2 = r1.get(v).text.split('"tdar">')[i][
                                                     :r1.get(v).text.split('"tdar">')[i].find('</td>')]
                                                zz += f'{x2}\n{x1}\n'
                                            i += 1

                                        bot.send_message(ch, zz)
                                    except:
                                        try:
                                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                            head = {
                                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                                'Accept-Encoding': 'gzip, deflate, br',
                                                'Accept-Language': 'en-US,en;q=0.9',
                                                'Connection': 'keep-alive',
                                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                                'Host': 'context.reverso.net',
                                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                                'Sec-Fetch-Dest': 'document',
                                                'Sec-Fetch-Mode': 'navigate',
                                                'Sec-Fetch-Site': 'same-origin',
                                                'Sec-Fetch-User': '?1',
                                                'Upgrade-Insecure-Requests': '1',
                                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                            }
                                            j = r1.get(url=url, headers=head)
                                            v = (j.text[
                                                 j.text.find(
                                                     '<button class="other-content" data-other="0" data-negative="') + len(
                                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                                     'Other translations</button>')])[:-2]
                                            c = (v.replace("-{", "").replace("}", "\n"))
                                            bot.send_message(msg.chat.id, f'* {c} *')
                                        except:
                                            pass
                                    head = {
                                        'accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'accept-encoding': 'gzip, deflate, br',
                                        'accept-language': 'en-US,en;q=0.9',
                                        'cookie': 'CGIC=IocBdGV4dC9odG1sLGFwcGxpY2F0aW9uL3hodG1sK3htbCxhcHBsaWNhdGlvbi94bWw7cT0wLjksaW1hZ2UvYXZpZixpbWFnZS93ZWJwLGltYWdlL2FwbmcsKi8qO3E9MC44LGFwcGxpY2F0aW9uL3NpZ25lZC1leGNoYW5nZTt2PWIzO3E9MC45; HSID=ACpCi--xKNWBuzc1V; SSID=AxTv5nnKkxKd-t24h; APISID=OvtsQIlYX38DaDAB/AejTrmzLNAmxJKW5L; SAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; __Secure-3PAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; CONSENT=YES+IQ.ar+201908; SEARCH_SAMESITE=CgQIoJAB; OTZ=5713657_44_44__44_; SID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGeJPLKTJluqScaUfUXXkqGg.; __Secure-3PSID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGofFKYNMh0FN74vl_7RXMLg.; ANID=AHWqTUl-D-kiT0ekElpQYvwaLkYpmfjEj0_vWVNyeQ9lR3vPGLYG1ea_vAtrh9jw; 1P_JAR=2020-11-26-23; NID=204=Nvh8sYQiAGl2xGl9Rivw7mwHK0oiirbBuBjelwbKer7spRkQS3Xh-9GYO5GR2WOwUG-GICaT4wz4vP8fWCpH8Zao3AEmb_XGvZ7CoLvMiTCbcSnMjY6GCGnFcwc9Ap2D0Iu1ltcyj4qzLl25BPgZf0KdrrpVYB_UyBb_LrJWA7A4dgrtBLEwtyBFzHe1XdYKS_yOA6QzeeHlNNtZm2YYpYUzvzevK2wxD2EMs9f9OOnq_es; DV=Q77RAgvgUI1KQP_f7OZE3e4dEcltYJdw4P5MyRCtgwEAAHAXk3Gqw6shngAAAAyZb21iQXUwRQAAAA; SIDCC=AJi4QfGRLCnLLJayzTS67kgRZvvsMJ8WJx4Z3VQIypirxBw8XSvY82_9ydx6j9Wr0SSSgWwkPw; __Secure-3PSIDCC=AJi4QfFg7EgRAZ_f0vxbunVqGgMst8xwoewWB6J0ZP6y8sihJ3OufFh34ZiNnfFg87jxfwhwgg',
                                        'referer': 'https://www.google.com/',
                                        'sec-fetch-dest': 'document',
                                        'sec-fetch-mode': 'navigate',
                                        'sec-fetch-site': 'same-origin',
                                        'sec-fetch-user': '?1',
                                        'upgrade-insecure-requests': '1',
                                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36',
                                        'x-client-data': 'CIy2yQEIpLbJAQjBtskBCKmdygEIlqzKAQisx8oBCPbHygEI6cjKAQi0y8oBCI3PygEI3NXKAQjul8sBCJGZywEImJrLARiKwcoB'
                                    }
                                    r = requests.get(
                                        f'https://www.google.com/search?q={d}+anatomy&source=lnms&tbm=isch&sa=X&ved=2ahUKEwjfv9Xg5YrvAhWHnxQKHe-2AfEQ_AUoAXoECBQQAw&biw=1821&bih=876#imgrc=3PIl4efjvoyhfM',
                                        headers=head)
                                    i = 0
                                    import shutil
                                    try:
                                        while i < 3:
                                            try:
                                                x = r.text.split('jpg"')[i].split('["http')[-1]
                                                v = f'http{x}jpg'

                                                with open(fr'koky{i}.jpg', 'wb') as mp3:
                                                    shutil.copyfileobj(r1.get(v, stream=True).raw, mp3)
                                                    bot.send_photo(chat_id=ch,
                                                                   photo=open(f'koky{i}.jpg', 'rb'))

                                            except:
                                                pass
                                            i += 1
                                    except:
                                        pass


                                else:
                                    url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                    head = {
                                        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'Accept-Encoding': 'gzip, deflate, br',
                                        'Accept-Language': 'en-US,en;q=0.9',
                                        'Connection': 'keep-alive',
                                        'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                        'Host': 'context.reverso.net',
                                        'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                        'Sec-Fetch-Dest': 'document',
                                        'Sec-Fetch-Mode': 'navigate',
                                        'Sec-Fetch-Site': 'same-origin',
                                        'Sec-Fetch-User': '?1',
                                        'Upgrade-Insecure-Requests': '1',
                                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                    }
                                    j = r1.get(url=url, headers=head)
                                    v = (j.text[
                                         j.text.find(
                                             '<button class="other-content" data-other="0" data-negative="') + len(
                                             '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                             'Other translations</button>')])[:-2]
                                    c = (v.replace("-{", "").replace("}", "\n"))
                                    bot.send_message(msg.chat.id, f'* {c} *')
                            except:
                                pass
                except:
                    pass
                if msg.reply_to_message:
                    pass
                elif 'لغاء' in msg.text:
                    if ch in g:
                        g.remove(ch)
                    else:
                        pass
                else:
                    try:
                        url = 'https://www.arabtran.com/gtranslate/'

                        head = {
                            'accept': '*/*',
                            'accept-encoding': 'gzip, deflate, br',
                            'accept-language': 'en-US,en;q=0.9',
                            'content-length': '31',
                            'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                            'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                            'origin': 'https://www.arabtran.com',
                            'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                            'sec-fetch-dest': 'empty',
                            'sec-fetch-mode': 'cors',
                            'sec-fetch-site': 'same-origin',
                            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                            'x-requested-with': 'XMLHttpRequest',
                        }

                        q = types.InlineKeyboardMarkup()
                        q1 = types.InlineKeyboardButton('اضغط هنا في حيال وجود مصطلح طبي 💟', callback_data='q1')
                        q.add(q1)
                        if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o', 'p',
                                         'q', 'r',
                                         's', 't', 'u', 'v', 'w', 'x', 'y', 'z'] or f'{w[1]}' in ['a', 'b', 'c', 'd',
                                                                                                  'e',
                                                                                                  'f', 'g',
                                                                                                  'h', 'i', 'j', 'k',
                                                                                                  'l',
                                                                                                  'm', 'n',
                                                                                                  'o', 'p', 'q', 'r',
                                                                                                  's',
                                                                                                  't', 'u',
                                                                                                  'v', 'w', 'x', 'y',
                                                                                                  'z'] or f'{w[2]}' in [
                            'a',
                            'b',
                            'c',
                            'd',
                            'e',
                            'f',
                            'g',
                            'h',
                            'i',
                            'j',
                            'k',
                            'l',
                            'm',
                            'n',
                            'o',
                            'p',
                            'q',
                            'r',
                            's',
                            't',
                            'u',
                            'v',
                            'w',
                            'x',
                            'y',
                            'z']:
                            data = {
                                'text': w,
                                'gfrom': 'en',
                                'gto': 'ar',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text
                            bot.send_message(msg.chat.id, x, reply_markup=q)

                        else:
                            data = {
                                'text': w,
                                'gfrom': 'ar',
                                'gto': 'en',
                                'key': 'ABC'
                            }
                            x = r1.post(url=url, data=data, ).text
                            bot.send_message(msg.chat.id, x, reply_markup=q)


                    except:
                        pass
            else:
                if 'تفعيل' in msg.text:
                    bot.reply_to(msg, 'تم تفعيل بوت الترجمة بنجاح 💚')
                    g.append(ch)
        else:
            if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                pass
            else:
                open('tran.txt', 'a').write(f'{msg.chat.id}\n')
            idd = msg.chat.id
            res = \
                r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                    'result'][
                    'status']
            if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
            else:
                try:
                    if msg.text == 'الاعضاء' and msg.chat.id in [idi, 1490464385]:
                        x = open('tran.txt', 'r').readlines()
                        bot.reply_to(msg, '{}'.format(len(x)))
                except:
                    pass
                try:
                    if msg.text == 'الاذاعة' and msg.chat.id == idi:
                        markup = types.ForceReply(selective=False)
                        bot.send_message(msg.chat.id, "ارسل اذاعتك", reply_markup=markup)
                except:
                    pass
                try:
                    if msg.reply_to_message.text == "ارسل اذاعتك":
                        try:
                            s = msg.text
                            x = open('tran.txt', 'r')
                            for i in x:
                                try:
                                    bot.send_message(i.replace('\n', ''), s)
                                except:
                                    pass
                        except:
                            pass
                except:
                    pass
                w = msg.text.lower()
                try:
                    if msg.reply_to_message.text == "ارسل المصطلح الطبي المراد ترجمته":
                        if w[0] == 'ا' or w[0] == 'أ' and w[1] == 'ل':
                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                            head = {
                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                'Accept-Encoding': 'gzip, deflate, br',
                                'Accept-Language': 'en-US,en;q=0.9',
                                'Connection': 'keep-alive',
                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                'Host': 'context.reverso.net',
                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                'Sec-Fetch-Dest': 'document',
                                'Sec-Fetch-Mode': 'navigate',
                                'Sec-Fetch-Site': 'same-origin',
                                'Sec-Fetch-User': '?1',
                                'Upgrade-Insecure-Requests': '1',
                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                            }
                            j = r1.get(url=url, headers=head)
                            v = (j.text[
                                 j.text.find('<button class="other-content" data-other="0" data-negative="') + len(
                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                     'Other translations</button>')])[:-2]
                            c = (v.replace("-{", "").replace("}", "\n"))
                            bot.send_message(msg.chat.id, f'* {c} *')

                        else:
                            try:
                                if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n',
                                                 'o',
                                                 'p',
                                                 'q',
                                                 'r',
                                                 's', 't', 'u', 'v', 'w', 'x', 'y', 'z']:
                                    d = w
                                    v = f'https://www.tbeeb.net/%D9%82%D8%A7%D9%85%D9%88%D8%B3-%D8%B7%D8%A8%D9%8A/search.php?q={d}+&dictionary=%D8%A8%D8%AD%D8%AB'
                                    try:
                                        i = 1
                                        zz = ''
                                        while i < len(r1.get(v).text.split('"tden">')):
                                            if i > 7:
                                                pass
                                            else:
                                                x1 = r1.get(v).text.split('"tden">')[i][
                                                     :r1.get(v).text.split('"tden">')[i].find('</td>')]
                                                x2 = r1.get(v).text.split('"tdar">')[i][
                                                     :r1.get(v).text.split('"tdar">')[i].find('</td>')]
                                                zz += f'{x2}\n{x1}\n'
                                            i += 1

                                        bot.send_message(ch, zz)
                                    except:
                                        try:
                                            url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                            head = {
                                                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                                'Accept-Encoding': 'gzip, deflate, br',
                                                'Accept-Language': 'en-US,en;q=0.9',
                                                'Connection': 'keep-alive',
                                                'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                                'Host': 'context.reverso.net',
                                                'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                                'Sec-Fetch-Dest': 'document',
                                                'Sec-Fetch-Mode': 'navigate',
                                                'Sec-Fetch-Site': 'same-origin',
                                                'Sec-Fetch-User': '?1',
                                                'Upgrade-Insecure-Requests': '1',
                                                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                            }
                                            j = r1.get(url=url, headers=head)
                                            v = (j.text[
                                                 j.text.find(
                                                     '<button class="other-content" data-other="0" data-negative="') + len(
                                                     '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                                     'Other translations</button>')])[:-2]
                                            c = (v.replace("-{", "").replace("}", "\n"))
                                            bot.send_message(msg.chat.id, f'* {c} *')
                                        except:
                                            pass
                                    head = {
                                        'accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'accept-encoding': 'gzip, deflate, br',
                                        'accept-language': 'en-US,en;q=0.9',
                                        'cookie': 'CGIC=IocBdGV4dC9odG1sLGFwcGxpY2F0aW9uL3hodG1sK3htbCxhcHBsaWNhdGlvbi94bWw7cT0wLjksaW1hZ2UvYXZpZixpbWFnZS93ZWJwLGltYWdlL2FwbmcsKi8qO3E9MC44LGFwcGxpY2F0aW9uL3NpZ25lZC1leGNoYW5nZTt2PWIzO3E9MC45; HSID=ACpCi--xKNWBuzc1V; SSID=AxTv5nnKkxKd-t24h; APISID=OvtsQIlYX38DaDAB/AejTrmzLNAmxJKW5L; SAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; __Secure-3PAPISID=VictV1NMEEknSmKC/Ap3tCyc9Rw-nNhkKN; CONSENT=YES+IQ.ar+201908; SEARCH_SAMESITE=CgQIoJAB; OTZ=5713657_44_44__44_; SID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGeJPLKTJluqScaUfUXXkqGg.; __Secure-3PSID=3gekKKpt2BDYs3HGt2dbGMJygRXv0S_uH30dI87nwwd9kPVGofFKYNMh0FN74vl_7RXMLg.; ANID=AHWqTUl-D-kiT0ekElpQYvwaLkYpmfjEj0_vWVNyeQ9lR3vPGLYG1ea_vAtrh9jw; 1P_JAR=2020-11-26-23; NID=204=Nvh8sYQiAGl2xGl9Rivw7mwHK0oiirbBuBjelwbKer7spRkQS3Xh-9GYO5GR2WOwUG-GICaT4wz4vP8fWCpH8Zao3AEmb_XGvZ7CoLvMiTCbcSnMjY6GCGnFcwc9Ap2D0Iu1ltcyj4qzLl25BPgZf0KdrrpVYB_UyBb_LrJWA7A4dgrtBLEwtyBFzHe1XdYKS_yOA6QzeeHlNNtZm2YYpYUzvzevK2wxD2EMs9f9OOnq_es; DV=Q77RAgvgUI1KQP_f7OZE3e4dEcltYJdw4P5MyRCtgwEAAHAXk3Gqw6shngAAAAyZb21iQXUwRQAAAA; SIDCC=AJi4QfGRLCnLLJayzTS67kgRZvvsMJ8WJx4Z3VQIypirxBw8XSvY82_9ydx6j9Wr0SSSgWwkPw; __Secure-3PSIDCC=AJi4QfFg7EgRAZ_f0vxbunVqGgMst8xwoewWB6J0ZP6y8sihJ3OufFh34ZiNnfFg87jxfwhwgg',
                                        'referer': 'https://www.google.com/',
                                        'sec-fetch-dest': 'document',
                                        'sec-fetch-mode': 'navigate',
                                        'sec-fetch-site': 'same-origin',
                                        'sec-fetch-user': '?1',
                                        'upgrade-insecure-requests': '1',
                                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36',
                                        'x-client-data': 'CIy2yQEIpLbJAQjBtskBCKmdygEIlqzKAQisx8oBCPbHygEI6cjKAQi0y8oBCI3PygEI3NXKAQjul8sBCJGZywEImJrLARiKwcoB'
                                    }
                                    r = requests.get(
                                        f'https://www.google.com/search?q={d}+anatomy&source=lnms&tbm=isch&sa=X&ved=2ahUKEwjfv9Xg5YrvAhWHnxQKHe-2AfEQ_AUoAXoECBQQAw&biw=1821&bih=876#imgrc=3PIl4efjvoyhfM',
                                        headers=head)
                                    i = 0
                                    import shutil
                                    try:
                                        while i < 3:
                                            try:
                                                x = r.text.split('jpg"')[i].split('["http')[-1]
                                                v = f'http{x}jpg'

                                                with open(fr'koky{i}.jpg', 'wb') as mp3:
                                                    shutil.copyfileobj(r1.get(v, stream=True).raw, mp3)
                                                    bot.send_photo(chat_id=ch,
                                                                   photo=open(f'koky{i}.jpg', 'rb'))

                                            except:
                                                pass
                                            i += 1
                                    except:
                                        pass


                                else:
                                    url = f'https://context.reverso.net/translation/arabic-english/{w}'
                                    head = {
                                        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.9',
                                        'Accept-Encoding': 'gzip, deflate, br',
                                        'Accept-Language': 'en-US,en;q=0.9',
                                        'Connection': 'keep-alive',
                                        'Cookie': 'didomi_token=eyJ1c2VyX2lkIjoiMTc1YWYxY2QtM2E4NS02NWNhLTlhNWYtNWMzMTMxYmY0OTQyIiwiY3JlYXRlZCI6IjIwMjAtMTEtMDlUMjI6MjU6MDMuODM3WiIsInVwZGF0ZWQiOiIyMDIwLTExLTA5VDIyOjI1OjAzLjgzN1oiLCJ2ZXJzaW9uIjpudWxsfQ==; __qca=P0-2069518181-1604960703577; _ga=GA1.2.1022455205.1604960704; _fbp=fb.1.1604960704785.1696161273; __gads=ID=6a415d82c09e7f2f:T=1604961014:S=ALNI_MYnjNMXo7nFRuuahoJIpqXamOQFvA; CTXTNODEID=bstweb12; experiment_context_N7gT3vKzX=0; JSESSIONID=kamXv4rTSPJG1mzG5D0fBsyY.bst-web12; context.lastpair=ar-en; history_entry=psycho]#[gastric]#[{w}; history_pair=en-ar]#[en-ar]#[ar-en; experiment_context_E3de3pqAZ=1; _gid=GA1.2.1607272187.1606403162; context.dapppromotion-count2=1; context.dapppromotion2=0; _gat_gtag_UA_2834324_41=1',
                                        'Host': 'context.reverso.net',
                                        'Referer': 'https://context.reverso.net/translation/arabic-english/%7Bw%7D',
                                        'Sec-Fetch-Dest': 'document',
                                        'Sec-Fetch-Mode': 'navigate',
                                        'Sec-Fetch-Site': 'same-origin',
                                        'Sec-Fetch-User': '?1',
                                        'Upgrade-Insecure-Requests': '1',
                                        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/86.0.4240.198 Safari/537.36'
                                    }
                                    j = r1.get(url=url, headers=head)
                                    v = (j.text[
                                         j.text.find(
                                             '<button class="other-content" data-other="0" data-negative="') + len(
                                             '<button class="other-content" data-other="0" data-negative="'):j.text.find(
                                             'Other translations</button>')])[:-2]
                                    c = (v.replace("-{", "").replace("}", "\n"))
                                    bot.send_message(msg.chat.id, f'* {c} *')
                            except:
                                pass


                except:
                    pass
                if msg.reply_to_message:
                    pass
                else:
                    try:
                        url = 'https://www.arabtran.com/gtranslate/'

                        head = {
                            'accept': '*/*',
                            'accept-encoding': 'gzip, deflate, br',
                            'accept-language': 'en-US,en;q=0.9',
                            'content-length': '31',
                            'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                            'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                            'origin': 'https://www.arabtran.com',
                            'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                            'sec-fetch-dest': 'empty',
                            'sec-fetch-mode': 'cors',
                            'sec-fetch-site': 'same-origin',
                            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                            'x-requested-with': 'XMLHttpRequest',
                        }

                        q = types.InlineKeyboardMarkup()
                        q1 = types.InlineKeyboardButton('اضغط هنا في حيال وجود مصطلح طبي 💟', callback_data='q1')
                        q.add(q1)
                        if f'{w[0]}' in ['a', 'b', 'c', 'd', 'e', 'f', 'g', 'h', 'i', 'j', 'k', 'l', 'm', 'n', 'o', 'p',
                                         'q',
                                         'r',
                                         's', 't', 'u', 'v', 'w', 'x', 'y', 'z'] or f'{w[1]}' in ['a', 'b', 'c', 'd',
                                                                                                  'e',
                                                                                                  'f',
                                                                                                  'g',
                                                                                                  'h', 'i', 'j', 'k',
                                                                                                  'l',
                                                                                                  'm',
                                                                                                  'n',
                                                                                                  'o', 'p', 'q', 'r',
                                                                                                  's',
                                                                                                  't',
                                                                                                  'u',
                                                                                                  'v', 'w', 'x', 'y',
                                                                                                  'z'] or f'{w[2]}' in [
                            'a',
                            'b',
                            'c',
                            'd',
                            'e',
                            'f',
                            'g',
                            'h',
                            'i',
                            'j',
                            'k',
                            'l',
                            'm',
                            'n',
                            'o',
                            'p',
                            'q',
                            'r',
                            's',
                            't',
                            'u',
                            'v',
                            'w',
                            'x',
                            'y',
                            'z']:
                            data = {
                                'text': w,
                                'gfrom': 'en',
                                'gto': 'ar',
                                'key': 'ABC'
                            }
                            x = GoogleTranslator(source='auto', target='ar').translate(
                                text=msg.text)
                            bot.send_message(msg.chat.id, x, reply_markup=q)

                        else:
                            data = {
                                'text': w,
                                'gfrom': 'ar',
                                'gto': 'en',
                                'key': 'ABC'
                            }
                            x = GoogleTranslator(source='auto', target='en').translate(
                                text=msg.text)

                            bot.send_message(msg.chat.id, x, reply_markup=q)

                    except:
                        pass
    except:
        pass


@bot.callback_query_handler(lambda call: True)
def any(call):
    if call.data == 'q1':
        markup = types.ForceReply(selective=False)
        bot.send_message(call.message.chat.id, "ارسل المصطلح الطبي المراد ترجمته", reply_markup=markup)


h = []


@bot.message_handler(content_types='document')
def an(msg):
    ch = msg.chat.id  #
    try:
        if ch in work:
            bot.send_message(ch, 'انتضر قليلا ليترجم الملف الذي ارسلته سابقا⚠️')
        else:
            work.append(ch)
            if msg.chat.type != "private":
                if ch in g:
                    if ch in c:
                        pass
                    else:
                        c.append(ch)

                    if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                        pass
                    else:
                        open('tran.txt', 'a').write(f'{msg.chat.id}\n')
                        bot.send_video(ch, 'BAACAgIAAxkBAAEBMbBgCBkFA1GYku_Q1wdBAAHTqU-ItsIAAhQMAAKHfEBIKhPv8pCI94oeBA')
                        bot.send_voice(ch, 'AwACAgIAAxkBAAEBMX9gCBK56Cj7teHoOhauUWlr7poIZQACEAwAAod8QEiJYsq2h5MWXR4E')
                    ch = msg.chat.id
                    user = msg.from_user.username
                    open(f'{user}.txt', 'w').write('Ahmed Koky (Q5QQQQ)\n\n')
                    open(f'{user}ar.txt', 'w', encoding='utf-8').write('')

                    try:
                        file_info = bot.get_file(msg.document.file_id)
                        downloaded_file = bot.download_file(file_info.file_path)
                        x = msg.document.file_name[-4:].lower()
                        if x == '.pdf' or x == '.PDF':
                            if ch in h:
                                bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                            else:
                                h.append(ch)
                                bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                                
                                with open('ahmed.pdf', 'wb') as new_file:
                                    new_file.write(downloaded_file)
                                with open('ahmed.pdf', 'rb') as f:
                                    # Create a PDF object
                                    pdf = PyPDF2.PdfFileReader(f)
                                    # Get the number of pages in the PDF
                                    num_pages = pdf.getNumPages()
                                    # Initialize a list to store the translated text
                                    translated_text = []
                                    # Iterate over each page in the PDF
                                    for page_num in range(num_pages):
                                        # Extract the text from the page
                                        page_text = pdf.getPage(page_num).extractText()

                                with fitz.open("ahmed.pdf") as doc:
                                    text = ""
                                    for page in doc:
                                        text += page.getText()
                                        open(f'{user}.txt', 'a').write(f'{text}\n')
                                    z = open(f'{user}.txt', 'r').read()
                                    if z.replace("\n", "").replace(" ", "") == '':
                                        pdf = wi(filename="ahmed.pdf", resolution=150)
                                        pdfImage = pdf.convert('jpeg')

                                        imageBlobs = []

                                        for img in pdfImage.sequence:
                                            imgPage = wi(image=img)
                                            imageBlobs.append(imgPage.make_blob('jpeg'))

                                        recognized_text = []

                                        for imgBlob in imageBlobs:
                                            im = Image.open(io.BytesIO(imgBlob))
                                            text = pytesseract.image_to_string(im, lang='eng')
                                            recognized_text.append(text)
                                        for k in recognized_text:
                                            open(f'{user}.txt', 'a').write(f'{k}\n')
                                    url = 'https://www.arabtran.com/gtranslate/'
                                    head = {
                                        'accept': '*/*',
                                        'accept-encoding': 'gzip, deflate, br',
                                        'accept-language': 'en-US,en;q=0.9',
                                        'content-length': '31',
                                        'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                                        'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                                        'origin': 'https://www.arabtran.com',
                                        'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                                        'sec-fetch-dest': 'empty',
                                        'sec-fetch-mode': 'cors',
                                        'sec-fetch-site': 'same-origin',
                                        'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                                        'x-requested-with': 'XMLHttpRequest',
                                    }
                                    z = open(f'{user}.txt', 'r').read()

                                    if len(z) > 5000:
                                        if ch in d:
                                            k = d.index(ch)
                                            try:
                                                d.remove(k)
                                                d.remove(k)
                                            except:
                                                pass
                                        else:
                                            d.append(ch)
                                            d.append(' ')

                                        t = ''
                                        e = d[d.index(ch) + 1]
                                        u = []
                                        p = 0
                                        while (p - 1) * 5000 < len(z):
                                            if p == 0:
                                                data = {
                                                    'text': z,
                                                    'gfrom': 'en',
                                                    'gto': 'ar',
                                                    'key': 'ABC'
                                                }
                                                x = r1.post(url=url, data=data, ).text
                                                t += x
                                                q = t.split('\n')
                                                q1 = z.split('\n')
                                                i = 0
                                                while i < len(q1):
                                                    try:
                                                        open('kok.txt', 'a', encoding='utf-8').write(
                                                            f'{q[i]}\n{q1[i]}\n')
                                                        e += f'{q1[i]}\n'
                                                        if i - len(q1) == -1:
                                                            d.insert(d.index(ch) + 1, e)
                                                    except:
                                                        pass
                                                    i += 1
                                            else:
                                                data = {
                                                    'text': z[z.find(d[d.index(ch) + 1]) + len(e):],
                                                    'gfrom': 'en',
                                                    'gto': 'ar',
                                                    'key': 'ABC'
                                                }
                                                x = r1.post(url=url, data=data, ).text
                                                q = x.split('\n')
                                                q1 = z[z.find(e) + len(e):].split('\n')
                                                i = 0
                                                while i < len(q1):
                                                    try:
                                                        open('kok.txt', 'a', encoding='utf-8').write(
                                                            f'{q[i]}\n{q1[i]}\n')
                                                        e += f'{q1[i]}\n'
                                                        if i - len(q1) == -1:
                                                            d.insert(d.index(ch) + 1, e)
                                                    except:
                                                        pass
                                                    i += 1
                                        p += 1
                                    else:
                                        data = {
                                            'text': z,
                                            'gfrom': 'en',
                                            'gto': 'ar',
                                            'key': 'ABC'
                                        }
                                        x = r1.post(url=url, data=data, ).text
                                        open(f'{user}ar.txt', 'a', encoding='utf-8').write(f"{x}\n")
                                    data = {
                                        'text': z,
                                        'gfrom': 'en',
                                        'gto': 'ar',
                                        'key': 'ABC'
                                    }
                                    x = open(f'{user}.txt', 'r').readlines()
                                    x1 = open(f'{user}ar.txt', 'r', encoding='utf-8').readlines()
                                    open('ahmed.txt', 'w').write('')
                                    i = 1
                                    while i < len(x):
                                        try:
                                            open('ahmed.txt', 'a', encoding='utf-8').write(f'{x1[i]}')
                                            open('ahmed.txt', 'a', encoding='utf-8').write(f'{x[i]}\n')
                                        except:
                                            pass
                                        i += 1
                                    bot.send_document(ch, open(f'ahmed.txt', 'r', encoding='utf-8'),
                                                      caption=msg.document.file_name)



                        else:
                            bot.send_message(ch, 'من فضلك ارسل الملفات بصيغ ال pdf ')

                    except:
                        pass
                    try:
                        h.remove(ch)
                    except:
                        pass
            else:
                if ch in c:
                    pass
                else:
                    c.append(ch)

                if f'{msg.chat.id}\n' in open('tran.txt', 'r'):
                    pass
                else:
                    open('tran.txt', 'a').write(f'{msg.chat.id}\n')
                    bot.send_video(ch, 'BAACAgIAAxkBAAEBMbBgCBkFA1GYku_Q1wdBAAHTqU-ItsIAAhQMAAKHfEBIKhPv8pCI94oeBA')
                    bot.send_voice(ch, 'AwACAgIAAxkBAAEBMX9gCBK56Cj7teHoOhauUWlr7poIZQACEAwAAod8QEiJYsq2h5MWXR4E')
                idd = msg.chat.id
                res = \
                    r1.get(f'https://api.telegram.org/bot{token}/getChatMember?chat_id=@akokybot&user_id={idd}').json()[
                        'result'][
                        'status']
                if res == 'left' and msg.chat.id not in [idi, 1490464385]:
                    bot.send_message(ch, 'يجب عليك الاشتراك في القناة اول و من ثم استخدم البوت \nلطفا💚\n@akokybot')
                else:
                    ch = msg.chat.id
                    user = msg.from_user.username
                    open(f'{user}.txt', 'w').write('')
                    open(f'{user}ar.txt', 'w', encoding='utf-8').write('')

                    try:
                        file_info = bot.get_file(msg.document.file_id)
                        downloaded_file = bot.download_file(file_info.file_path)
                        x = msg.document.file_name[-4:]
                        if x == '.pdf' or x == '.PDF':

                            if ch in h:

                                bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                            else:
                                h.append(ch)
                                bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                                with open('ahmed.pdf', 'wb') as new_file:
                                    new_file.write(downloaded_file)
                                v = 0
                                document = Document()    
                                with open('ahmed.pdf', 'rb') as f:
                                    # Create a PDF object
                                    pdf = PyPDF2.PdfFileReader(f)
                                    # Get the number of pages in the PDF
                                    num_pages = pdf.getNumPages()
                                    # Initialize a list to store the translated text
                                    translated_text = []
                                    # Iterate over each page in the PDF
                                    txt= ''
                                    for page_num in range(num_pages):
                                        # Extract the text from the page
                                        page_text = pdf.getPage(page_num).extractText()
                                        page_text += '-----------------------'
                                        if len(page_text) > 5000:
                                            trans = ''
                                            i = len(page_text)//5000

                                            ii = 0
                                            iii = 0
                                            while ii < i + 1:

                                                tran = GoogleTranslator(source='en', target='ar').translate(
                                                    text=text[iii:4999 + iii])


                                                ii += 1
                                                iii += 5000
                                            trans += tran
                                        else:

                                            trans = GoogleTranslator(source='en', target='ar').translate(
                                                text=page_text)
                                        
                                        
                                
                                        document.add_paragraph(f'{trans}')

                                        open(f'{user}.txt', 'a', encoding='utf-8').write(text)
                                        open(f'{user}ar.txt', 'a', encoding='utf-8').write(trans)
                                        zx = 0
                                        zx1 = text.split('\n')
                                        zx2 = trans.split('\n')
                                        while len(text.split('\n')) > zx:
                                            try:
                                                strr += f"{zx1[zx]}\n{zx2[zx]} \n"
                                            except:
                                                pass
                                            zx += 1

                                        open(f'ahmed2.txt', 'a', encoding='utf-8').write(strr)
                                        print(strr)
                                    try:
                                        os.remove(f"{ch}.png")

                                    except:
                                        pass
                                    document.save('Ahmed2.docx')
                                    bot.send_document(ch, open(f'Ahmed2.docx', 'rb'), caption=f'{msg.document.file_name}\nترجمة من دون تغير الملف ')
                                    x1 = open(f'{user}ar.txt', 'r', encoding='utf-8').read()
                                    z = open(f'{user}.txt', 'r', encoding='utf-8').read()
                                    open('ahmed.txt', 'w', encoding='utf-8').write(x1)
                                    bot.send_document(ch, open(f'ahmed.txt', 'r', encoding='utf-8'),
                                                      caption=f'{msg.document.file_name}\nترجمة عربي فقط')
                                    x2 = open(f'{user}.txt', 'r', encoding='utf-8').read()
                                    cc = 0
                                    u = ''
                                    while cc < len(x2.split('|#|')):
                                        try:
                                            u += f"{z.split('|#|')[cc]}\n{x1.split('| # |')[cc]}\n"
                                        except:
                                            pass
                                        cc += 1

                                    document = Document()
                                    paragraph = document.add_paragraph(str(u))
                                    paragraph.style = document.styles.add_style('', WD_STYLE_TYPE.PARAGRAPH)
                                    font = paragraph.style.font
                                    font.size = Pt(16)
                                    font.bold = True
                                    font.color.rgb = RGBColor(0, 0, 0)
                                    document.save('Ahmed1.docx')
                                    bot.send_document(ch, open(f'Ahmed1.docx', 'rb'),
                                                      caption=f'{msg.document.file_name}\nترجمة صفحة صفحة')
                                    cc = 0
                                    u = open(f'ahmed2.txt', 'r', encoding='utf-8').read()
                                    print(u)

                                    document = Document()
                                    paragraph = document.add_paragraph(str(u))
                                    paragraph.style = document.styles.add_style('', WD_STYLE_TYPE.PARAGRAPH)
                                    font = paragraph.style.font
                                    font.size = Pt(16)
                                    font.bold = True
                                    font.color.rgb = RGBColor(0, 0, 0)
                                    document.save('Ahmed1.docx')
                                    bot.send_document(ch, open(f'Ahmed1.docx', 'rb'),
                                                      caption=f'{msg.document.file_name}\nترجمة سطرية')

                        else:
                            if ch in h:
                                bot.send_message(ch, 'انتضر قليلا من فضلك لكي يتم التحميل 💚')
                            else:
                                h.append(ch)
                                bot.send_message(ch, 'حسنا, سيتم التحميل الرجاء انتظار دقيقة واحدة 💚')
                            pp = []
                            open(f'{user}.txt', 'w').write('1')

                            x = (msg.document.file_name).split('.')[-1]

                            if x == 'pptx' or x == 'PPT':
                                with open('Ahmed.pptx', 'wb') as new_file:
                                    new_file.write(downloaded_file)

                                for eachfile in glob.glob("*.pptx"):

                                    prs = Presentation(eachfile)
                                    open(f'{user}.txt', 'a').write(
                                        'Done translation by Ahmed bot \n@any_kokybot\n@Q5QQQQ\n')
                                    open(f'{user}.txt', 'a').write('-----------------------\n')
                                    u = ''
                                    i = 0
                                    for slide in prs.slides:

                                        for shape in slide.shapes:
                                            if hasattr(shape, "text"):
                                                u += shape.text
                                        u += '---------'
                                        aa = u.split('---------')

                                        trans = GoogleTranslator(source='en', target='ar').translate(
                                                        text=str(aa[i]))

                                        i += 1


                                        open(f'{user}.txt', 'a', encoding='utf-8').write(f'{u}\n----------------\n')
                                        open(f'{user}ar.txt', 'a', encoding='utf-8').write(f'{trans}\n----------------\n')


                            elif x == 'docx' or x == 'DOCx':
                                u = ''
                                with open('Ahmed.docx', 'wb') as new_file:
                                    new_file.write(downloaded_file)
                                document = Document('Ahmed.docx')
                                open(f'{user}.txt', 'a', encoding='utf-8').write(
                                    'Done translation by Ahmed bot \n@any_kokybot\n@Q5QQQQ\n')
                                open(f'{user}.txt', 'a', encoding='utf-8').write('-----------------------\n')
                                for p in document.paragraphs:
                                    u = p.text
                                    u += '\n'
                                    u += '---------'
                                    if len(u) > 5000:
                                        print('ok')
                                    trans = GoogleTranslator(source='en', target='ar').translate(
                                        text=str(u))
                                    open(f'{user}.txt', 'a', encoding='utf-8').write(f'{u}\n----------------\n')
                                    open(f'{user}ar.txt', 'a', encoding='utf-8').write(
                                        f'{trans}\n')


                            else:
                                bot.send_message(ch, 'عذرا لم يتم ترجمة ملفك لانه ليس✳️ \npdf or pptx or docx\n')


                            url = 'https://www.arabtran.com/gtranslate/'
                            head = {
                                'accept': '*/*',
                                'accept-encoding': 'gzip, deflate, br',
                                'accept-language': 'en-US,en;q=0.9',
                                'content-length': '31',
                                'content-type': 'application/x-www-form-urlencoded; charset=UTF-8',
                                'cookie': '_ga=GA1.2.189101859.1607603223; _gid=GA1.2.154684385.1607603223; __gads=ID=ba2b5a0fe52a4c1b-22c3694088a60095:T=1607603223:RT=1607603223:S=ALNI_MbdZ0H2MeiuATaH2_vAl4hrUnPm8Q; _gat=1; sc_is_visitor_unique=rx12068393.1607603329.CA2EB2FDB25A4FB19C08D5DA0B092EBC.1.1.1.1.1.1.1.1.1',
                                'origin': 'https://www.arabtran.com',
                                'referer': 'https://www.arabtran.com/tarjamat_anjilizi_earabiun/',
                                'sec-fetch-dest': 'empty',
                                'sec-fetch-mode': 'cors',
                                'sec-fetch-site': 'same-origin',
                                'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/87.0.4280.88 Safari/537.36',
                                'x-requested-with': 'XMLHttpRequest',
                            }
                            z = open(f'{user}.txt', 'r', encoding='utf-8').read()

                            x = open(f'{user}ar.txt', 'r', encoding='utf-8').read().split('\n')
                            xx = z.split('\n')
                            i = 0
                            open('ahmed.txt', 'w', encoding='utf-8').write('')
                            while i < len(x):
                                open(f'ahmed.txt', 'a', encoding='utf-8').write(f"{x[i]}\n{xx[i]}\n")
                                i += 1
                            open('ahmed.txt', 'w', encoding='utf-8').write(
                                open(f'{user}ar.txt', 'r', encoding='utf-8').read())
                            bot.send_document(ch, open(f'ahmed.txt', 'r', encoding='utf-8'),
                                              caption=msg.document.file_name)
                            document = Document()

                            paragraph = document.add_paragraph(str(open(f'ahmed.txt', 'r', encoding='utf-8').read()))

                            paragraph.style = document.styles.add_style('', WD_STYLE_TYPE.PARAGRAPH)
                            font = paragraph.style.font
                            font.size = Pt(15)
                            font.bold = True
                            font.color.rgb = RGBColor(0, 0, 100)
                            document.save('Ahmed1.docx')
                            bot.send_document(ch, open(f'Ahmed1.docx', 'rb'),
                                              caption=msg.document.file_name)


                    except:
                        pass
                    try:
                        h.remove(ch)
                    except:
                        pass
            q = types.InlineKeyboardMarkup()
            a9 = types.InlineKeyboardButton("المطور 💚", url="https://t.me/Q5QQQQ")
            a10 = types.InlineKeyboardButton("قناه الاعلانات↗️", url="https://t.me/akokybot")
            q.add(a9)
            q.add(a10)
            try:
                os.remove(f'{user}.txt')
                os.remove(f'{user}ar.txt')
            except:
                pass
            bot.send_message(ch, 'الإعلانات ⬇,️', reply_markup=q)
        try:
            work.remove(ch)
        except:
            pass
    except:
        pass


i = 0
while i == 0:
    try:
        bot.polling(none_stop=True)
        i += 1
    except:
        pass
